"""
专业 PPT 排版引擎 v3 —— 智能多布局

每页内容会根据页类型自动选择最佳布局：
  cover   → 非对称封面（大色块+几何装饰）
  content → 智能选择：标准要点 / 两栏分列 / 居中大标题 / 步骤流程
  cards   → 圆角卡片网格
  chapter → 全幅色块章节分隔
  summary → 深色底总结页
"""

from __future__ import annotations

from abc import ABC
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement

from models.enums import AnimationType, TransitionType, ChartType


# 像素到 EMU 转换（1 像素 = 9525 EMU）
def Px(val):
    """将像素值转换为 EMU（PPT 内部单位）"""
    return Emu(val * 9525)


class BaseTemplate(ABC):

    # ── 元信息 ──
    id: str = ""
    name: str = ""
    description: str = ""
    thumbnail: str = ""

    # ── 配色 ──
    primary = (45, 100, 200)
    secondary = (80, 150, 240)
    accent = (255, 140, 50)
    bg = (255, 255, 255)
    dark_bg = (28, 35, 50)
    text = (33, 37, 41)
    text_light = (108, 117, 125)
    text_on_dark = (235, 240, 248)
    light_gray = (245, 248, 252)
    border = (222, 226, 230)

    font_cn = "Microsoft YaHei"
    font_en = "Arial"

    default_animation = AnimationType.FADE
    default_transition = TransitionType.FADE

    SW = Inches(13.333)
    SH = Inches(7.5)

    # ══════════════════════════════════════
    # 基础绘图
    # ══════════════════════════════════════

    def apply_theme(self, prs):
        prs.slide_width = self.SW
        prs.slide_height = self.SH

    def _slide(self, prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _bg(self, slide, c=None):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*(c or self.bg))

    def _rect(self, slide, l, t, w, h, color):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*color); s.line.fill.background()
        return s

    def _round_rect(self, slide, l, t, w, h, color):
        s = slide.shapes.add_shape(5, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*color); s.line.fill.background()
        return s

    def _circle(self, slide, l, t, d, color):
        s = slide.shapes.add_shape(9, l, t, d, d)
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(*color); s.line.fill.background()
        return s

    def _hline(self, slide, l, t, w, color, h=None):
        return self._rect(slide, l, t, w, h or Pt(1.5), color)

    def _txt(self, slide, l, t, w, h, text, sz=14, bold=False, color=None, align="l", font=None):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(sz); p.font.bold = bold
        p.font.color.rgb = RGBColor(*(color or self.text))
        p.font.name = font or self.font_cn
        if align == "c": p.alignment = PP_ALIGN.CENTER
        elif align == "r": p.alignment = PP_ALIGN.RIGHT
        return tb

    def _bullets(self, slide, l, t, w, h, items, sz=15, color=None, bullet="•", spacing=Pt(32)):
        if not items: return None
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        c = RGBColor(*(color or self.text))
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{bullet}  {item}"
            p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = self.font_cn
            p.space_after = Pt(10); p.line_spacing = spacing
        return tb

    # ══════════════════════════════════════
    # 封面 —— 非对称
    # ══════════════════════════════════════

    def layout_cover(self, prs, title, subtitle=""):
        slide = self._slide(prs)

        # 右半区大色块
        block_w = Inches(5.5)
        self._rect(slide, self.SW - block_w, Inches(0), block_w, self.SH, self.primary)

        # 色块内装饰圆
        self._circle(slide, self.SW - Inches(1.5), Inches(5.5), Inches(2.2), self.secondary)

        # 左上角小三角装饰
        self._rect(slide, Inches(0.6), Inches(0.6), Inches(1.2), Inches(0.06), self.accent)

        # 标题（左侧 60%）
        self._txt(slide, Inches(0.8), Inches(1.6), Inches(7.0), Inches(2.2),
                  title, sz=44, bold=True, color=self.text)
        self._hline(slide, Inches(0.8), Inches(4.0), Inches(2.8), self.primary, Pt(4))

        if subtitle:
            self._txt(slide, Inches(0.8), Inches(4.3), Inches(7.0), Inches(1.0),
                      subtitle, sz=17, color=self.text_light)

        self._txt(slide, Inches(0.8), Inches(6.6), Inches(5.0), Inches(0.4),
                  f"{self.name}  ·  PPT 智造", sz=10, color=self.text_light)
        return slide

    # ══════════════════════════════════════
    # 内容页 —— 智能布局选择
    # ══════════════════════════════════════

    def layout_content(self, prs, title, points):
        """根据要点数量和长度自动选择最佳布局"""
        n = len(points) if points else 0
        if n == 0:
            return self._layout_title_only(prs, title)
        elif n <= 2:
            return self._layout_two_column(prs, title, points)
        elif n <= 4:
            return self._layout_standard(prs, title, points)
        elif n <= 6:
            return self._layout_grid2(prs, title, points)
        else:
            return self._layout_grid3(prs, title, points)

    # ── 子布局 1：居中大标题（无要点时）──

    def _layout_title_only(self, prs, title):
        slide = self._slide(prs)
        self._bg(slide)
        self._hline(slide, Inches(0), Inches(0), self.SW, self.primary, Pt(5))
        self._txt(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.0),
                  title, sz=36, bold=True, color=self.text, align="c")
        self._hline(slide, Inches(5.0), Inches(5.2), Inches(3.3), self.secondary, Pt(3))
        return slide

    # ── 子布局 2：标准要点（3-4 项）──

    def _layout_standard(self, prs, title, points):
        slide = self._slide(prs)
        self._bg(slide)
        self._hline(slide, Inches(0), Inches(0), self.SW, self.primary, Pt(5))
        self._vline = lambda slide, l, t, w, h, color: self._rect(slide, l, t, w, h, color)
        self._vline(slide, Inches(0.55), Inches(0.6), Pt(6), Pt(50), self.primary)
        self._txt(slide, Inches(0.9), Inches(0.4), Inches(11.5), Inches(0.7),
                  title, sz=28, bold=True, color=self.primary)
        self._hline(slide, Inches(0.9), Inches(1.3), Inches(11.5), self.border, Pt(1.5))
        self._bullets(slide, Inches(1.2), Inches(1.7), Inches(11.0), Inches(5.2),
                      points, sz=16, bullet="▸", spacing=Pt(36))
        self._txt(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
                  self.name, sz=8, color=self.text_light, align="r", font=self.font_en)
        return slide

    # ── 子布局 3：两栏分列（1-2 项）──

    def _layout_two_column(self, prs, title, points):
        slide = self._slide(prs)
        self._bg(slide)
        self._hline(slide, Inches(0), Inches(0), self.SW, self.primary, Pt(5))
        self._txt(slide, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.7),
                  title, sz=28, bold=True, color=self.primary)
        self._hline(slide, Inches(0.8), Inches(1.2), Inches(11.5), self.border, Pt(1.5))

        col_w = Inches(5.5)
        col_h = Inches(5.0)
        for i, pt in enumerate(points):
            x = Inches(0.8) + (col_w + Inches(0.5)) * i
            # 栏背景
            self._round_rect(slide, x, Inches(1.6), col_w, col_h, self.light_gray)
            # 编号圆圈
            cx = x + col_w // 2 - Inches(0.3)
            self._circle(slide, cx, Inches(2.0), Inches(0.6), self.primary)
            self._txt(slide, cx, Inches(2.05), Inches(0.6), Inches(0.5),
                      str(i + 1), sz=22, bold=True, color=(255, 255, 255), align="c", font=self.font_en)
            # 内容
            self._txt(slide, x + Inches(0.5), Inches(2.9), col_w - Inches(1.0), Inches(3.3),
                      pt, sz=15, color=self.text)
        return slide

    # ── 子布局 4：2 列网格（5-6 项）──

    def _layout_grid2(self, prs, title, points):
        return self._build_grid(prs, title, points, cols=2)

    # ── 子布局 5：3 列网格（7+ 项）──

    def _layout_grid3(self, prs, title, points):
        return self._build_grid(prs, title, points, cols=3)

    def _build_grid(self, prs, title, points, cols):
        slide = self._slide(prs)
        self._bg(slide)
        self._hline(slide, Inches(0), Inches(0), self.SW, self.primary, Pt(5))
        self._txt(slide, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.7),
                  title, sz=28, bold=True, color=self.primary)
        self._hline(slide, Inches(0.8), Inches(1.2), Inches(11.5), self.border, Pt(1.5))

        gap = Inches(0.25)
        margin = Inches(0.6)
        bottom_pad = Inches(0.3)
        top = Inches(1.55)
        rows = (len(points) + cols - 1) // cols
        card_w = int((self.SW - margin * 2 - gap * (cols - 1)) / cols)
        usable_h = int(self.SH - top - gap * (rows - 1) - bottom_pad)
        card_h = int(usable_h / rows)

        for idx, pt in enumerate(points):
            row = idx // cols; col = idx % cols
            x = margin + col * (card_w + gap)
            y = top + row * (card_h + gap)

            self._round_rect(slide, x, y, card_w, card_h, self.light_gray)
            # 编号圆圈
            cx = x + card_w // 2 - Inches(0.25)
            cy = y + Inches(0.25)
            self._circle(slide, cx, cy, Inches(0.5), self.primary)
            self._txt(slide, cx, cy + Pt(2), Inches(0.5), Inches(0.4),
                      str(idx + 1), sz=16, bold=True, color=(255, 255, 255), align="c", font=self.font_en)
            # 文本
            text_w = card_w - Inches(0.6)
            self._txt(slide, x + Inches(0.3), y + Inches(1.0), text_w, card_h - Inches(1.3),
                      pt, sz=13, color=self.text)

        return slide

    # ══════════════════════════════════════
    # 卡片页（显式指定）
    # ══════════════════════════════════════

    def layout_cards(self, prs, title, points):
        return self._build_grid(prs, title, points,
                                cols=3 if len(points) >= 6 else (2 if len(points) >= 4 else 1))

    # ══════════════════════════════════════
    # 章节分隔
    # ══════════════════════════════════════

    def layout_chapter(self, prs, title, points):
        slide = self._slide(prs)
        self._bg(slide, self.primary)
        self._circle(slide, Inches(1.0), Inches(2.0), Inches(2.8), self.secondary)
        self._txt(slide, Inches(1.0), Inches(3.0), Inches(2.8), Inches(0.5),
                  "●", sz=50, bold=True, color=(255, 255, 255), align="c")
        self._txt(slide, Inches(4.2), Inches(2.2), Inches(8.0), Inches(1.5),
                  title, sz=40, bold=True, color=(255, 255, 255))
        self._hline(slide, Inches(4.2), Inches(3.9), Inches(3.5), (255, 255, 255), Pt(2))
        if points:
            self._txt(slide, Inches(4.2), Inches(4.3), Inches(8.0), Inches(2.0),
                      "  ·  ".join(points[:3]), sz=14, color=(200, 215, 235))
        return slide

    # ══════════════════════════════════════
    # 总结页
    # ══════════════════════════════════════

    def layout_summary(self, prs, title, points):
        slide = self._slide(prs)
        self._bg(slide, self.dark_bg)
        self._txt(slide, Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0),
                  title, sz=38, bold=True, color=(255, 255, 255), align="c")
        self._hline(slide, Inches(5.0), Inches(2.0), Inches(3.3), self.secondary, Pt(3))
        if points:
            self._bullets(slide, Inches(2.0), Inches(2.8), Inches(9.3), Inches(3.5),
                          points, sz=18, color=(225, 232, 242), bullet="✓", spacing=Pt(42))
        self._txt(slide, Inches(4.0), Inches(6.5), Inches(5.3), Inches(0.5),
                  "感谢聆听  ·  Thank You", sz=16, color=(170, 185, 205), align="c")
        return slide

    # ══════════════════════════════════════
    # 图表页
    # ════════════════════════════════════

    def layout_chart(self, prs, title, points, chart_data):
        """图表页面布局 - 标题 + 图表 + 可选要点说明"""
        slide = self._slide(prs)
        self._bg(slide)
        self._hline(slide, Inches(0), Inches(0), self.SW, self.primary, Pt(5))
        
        # 标题
        self._txt(slide, Inches(0.8), Inches(0.35), Inches(11.5), Inches(0.7),
                  title, sz=28, bold=True, color=self.primary)
        self._hline(slide, Inches(0.8), Inches(1.2), Inches(11.5), self.border, Pt(1.5))
        
        # 添加图表
        chart_title = chart_data.title if chart_data else ""
        self._add_chart(slide, chart_data, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), chart_title)
        
        # 如果有要点，在图表下方显示
        if points and len(points) > 0:
            # 预留较小的空间给要点
            bullet_top = Inches(6.6)
            bullet_h = Inches(0.6) * min(len(points), 3)
            self._bullets(slide, Inches(1.0), bullet_top, Inches(11.0), bullet_h,
                          points[:3], sz=13, color=self.text_light, bullet="•", spacing=Pt(24))
        
        return slide

    def _add_chart(self, slide, chart_data, left, top, width, height, title=""):
        """向幻灯片添加图表"""
        if not chart_data:
            return None
        
        # 处理 Pydantic 模型或字典
        if hasattr(chart_data, 'type'):
            # Pydantic 模型
            chart_type = chart_data.type
            categories = chart_data.categories
            series_data = chart_data.series
        else:
            # 字典
            chart_type = chart_data.get("type", "column")
            categories = chart_data.get("categories", [])
            series_data = chart_data.get("series", [])
        
        if not categories or not series_data:
            return None
        
        # 根据图表类型创建图表
        if chart_type in ("pie", "doughnut"):
            # 饼图/环形图 - 使用 CategoryChartData 作为数据源
            chart_data_obj = CategoryChartData()
            
            # 使用第一个系列的数据
            if series_data:
                first_series = series_data[0]
                # 处理字典类型（Pydantic 模型会自动解析为 dict）
                if isinstance(first_series, dict):
                    values = first_series.get("values", [])
                else:
                    values = getattr(first_series, 'values', [])
                
                # 添加虚拟分类
                chart_data_obj.categories = categories
                
                # 添加数据系列
                valid_values = []
                for v in values:
                    try:
                        valid_values.append(float(v))
                    except (ValueError, TypeError):
                        valid_values.append(0)
                
                if valid_values:
                    chart_data_obj.add_series("数据", valid_values)
            
            if len(chart_data_obj.categories) > 0:
                # 创建饼图
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.PIE if chart_type == "pie" else XL_CHART_TYPE.DOUGHNUT,
                    left, top, width, height, chart_data_obj
                ).chart
                chart.has_title = True
                chart.chart_title.text_frame.text = title
                if chart.legend:
                    chart.legend.position = XL_LEGEND_POSITION.RIGHT
                self._style_chart(chart)
        
        else:
            # 柱状图/条形图/折线图/面积图
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = categories
            
            for series in series_data:
                # 处理字典类型（Pydantic 模型会自动解析为 dict）
                if isinstance(series, dict):
                    series_name = series.get("name", "")
                    values = series.get("values", [])
                else:
                    # Pydantic 模型或其他对象
                    series_name = getattr(series, 'name', "")
                    values = getattr(series, 'values', [])
                
                # 过滤有效数值
                valid_values = []
                for v in values:
                    try:
                        valid_values.append(float(v))
                    except (ValueError, TypeError):
                        valid_values.append(0)
                
                chart_data_obj.add_series(series_name, valid_values)
            
            # 选择图表类型
            xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED
            if chart_type == "bar":
                xl_type = XL_CHART_TYPE.BAR_CLUSTERED
            elif chart_type == "line":
                xl_type = XL_CHART_TYPE.LINE_MARKERS
            elif chart_type == "area":
                xl_type = XL_CHART_TYPE.AREA
            
            chart = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data_obj).chart
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            if chart.legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            self._style_chart(chart)
        
        return chart

    def _style_chart(self, chart):
        """美化图表样式"""
        try:
            # 设置标题样式
            if chart.has_title:
                title_frame = chart.chart_title.text_frame
                for paragraph in title_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(*self.text)
            
            # 设置图例样式
            if chart.legend:
                for entry in chart.legend.entries:
                    entry.font.size = Pt(10)
                    entry.font.color.rgb = RGBColor(*self.text_light)
            
            # 设置坐标轴标签样式
            if hasattr(chart, 'category_axis'):
                chart.category_axis.tick_labels.font.size = Pt(10)
                chart.category_axis.tick_labels.font.color.rgb = RGBColor(*self.text_light)
            
            if hasattr(chart, 'value_axis'):
                chart.value_axis.tick_labels.font.size = Pt(10)
                chart.value_axis.tick_labels.font.color.rgb = RGBColor(*self.text_light)
                chart.value_axis.has_major_gridlines = True
            
            # 设置数据系列颜色
            for i, series in enumerate(chart.series):
                colors = [self.primary, self.secondary, self.accent, (100, 200, 100), (200, 100, 150)]
                color = colors[i % len(colors)]
                for point in series.points:
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = RGBColor(*color)
        
        except Exception:
            # 忽略样式设置错误
            pass

    # ══════════════════════════════════════

    def get_layout_method(self, page_type: str):
        return {
            "cover": self.layout_cover,
            "content": self.layout_content,
            "chapter": self.layout_chapter,
            "summary": self.layout_summary,
            "cards": self.layout_cards,
            "chart": self.layout_chart,
        }.get(page_type, self.layout_content)

    def set_transition(self, slide, trans_type=TransitionType.FADE):
        try:
            el = slide.element.find(qn("p:transition"))
            if el is None:
                el = slide.element.makeelement(qn("p:transition"), {})
                slide.element.append(el)
            for c in list(el): el.remove(c)
            m = {TransitionType.FADE: "fade", TransitionType.PUSH: "push",
                 TransitionType.WIPE: "wipe", TransitionType.ZOOM: "zoom"}
            if m.get(trans_type, "fade"): el.set(qn("p:advTm"), "500")
        except Exception: pass
